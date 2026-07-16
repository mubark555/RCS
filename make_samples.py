#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
إنشاء ملفات تجريبية للتجربة:
  * template.pptx : قالب فيه مكان محجوز {{name}}
  * names.xlsx    : ملف Excel فيه قائمة أسماء

للتجربة السريعة:
  python make_samples.py
  python name_inserter.py --pptx template.pptx --excel names.xlsx --output result.pptx --has-header
"""

from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from openpyxl import Workbook


def make_template(path="template.pptx"):
    prs = Presentation()
    prs.slide_width = Pt(720)
    prs.slide_height = Pt(540)

    slide = prs.slides.add_slide(prs.slide_layouts[6])  # شريحة فارغة

    # عنوان الشهادة
    box = slide.shapes.add_textbox(Pt(60), Pt(120), Pt(600), Pt(80))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "شهادة شكر وتقدير"
    run.font.size = Pt(40)
    run.font.bold = True

    # المكان المحجوز للاسم
    box2 = slide.shapes.add_textbox(Pt(60), Pt(240), Pt(600), Pt(80))
    tf2 = box2.text_frame
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = "{{name}}"
    run2.font.size = Pt(36)
    run2.font.bold = True

    prs.save(path)
    print(f"تم إنشاء القالب: {path}")


def make_names(path="names.xlsx"):
    wb = Workbook()
    ws = wb.active
    ws.title = "Names"
    ws["A1"] = "الاسم"  # عنوان العمود
    for i, name in enumerate(
        ["محمد العتيبي", "سارة القحطاني", "عبدالله الشهري", "نورة الزهراني", "فيصل الدوسري"],
        start=2,
    ):
        ws[f"A{i}"] = name
    wb.save(path)
    print(f"تم إنشاء ملف الأسماء: {path}")


if __name__ == "__main__":
    make_template()
    make_names()
