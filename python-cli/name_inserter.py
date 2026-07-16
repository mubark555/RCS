#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
name_inserter — إدراج الأسماء من ملف Excel داخل قالب PowerPoint.

الفكرة ببساطة:
  * عندك قالب PowerPoint فيه مكان محجوز للاسم (Placeholder) مثل: {{name}}
  * وعندك ملف Excel فيه قائمة الأسماء (عمود واحد).
  * الأداة تقرأ الأسماء، وتنسخ شريحة القالب لكل اسم، وتستبدل المكان المحجوز
    بالاسم، ثم تحفظ الناتج في ملف PowerPoint جديد (شريحة لكل اسم).

مثال:
  python name_inserter.py --pptx template.pptx --excel names.xlsx --output result.pptx

خيارات مفيدة:
  --placeholder "{{name}}"   النص المحجوز داخل القالب (الافتراضي {{name}})
  --column A                 عمود الأسماء في Excel (حرف أو رقم، الافتراضي أول عمود)
  --sheet "Sheet1"           اسم الورقة (الافتراضي أول ورقة)
  --has-header               تجاهُل الصف الأول (عنوان العمود)
  --slide 1                  رقم شريحة القالب التي تُنسخ (الافتراضي 1)
  --split                    إنشاء ملف مستقل لكل اسم بدلًا من ملف واحد
"""

import argparse
import copy
import os
import sys

try:
    from pptx import Presentation
except ImportError:
    sys.exit("خطأ: مكتبة python-pptx غير مثبتة. نفّذ: pip install python-pptx openpyxl")

try:
    from openpyxl import load_workbook
    from openpyxl.utils import column_index_from_string
except ImportError:
    sys.exit("خطأ: مكتبة openpyxl غير مثبتة. نفّذ: pip install python-pptx openpyxl")


# ---------------------------------------------------------------------------
# قراءة الأسماء من Excel
# ---------------------------------------------------------------------------
def read_names(excel_path, column=None, sheet=None, has_header=False):
    """يرجع قائمة بالأسماء (نصوص غير فارغة) من ملف Excel."""
    if not os.path.exists(excel_path):
        sys.exit(f"خطأ: ملف Excel غير موجود: {excel_path}")

    wb = load_workbook(excel_path, data_only=True, read_only=True)
    ws = wb[sheet] if sheet else wb.active

    # تحديد رقم العمود (الافتراضي: أول عمود = 1)
    if column is None:
        col_idx = 1
    elif str(column).isdigit():
        col_idx = int(column)
    else:
        col_idx = column_index_from_string(str(column).strip().upper())

    names = []
    for i, row in enumerate(ws.iter_rows(min_col=col_idx, max_col=col_idx, values_only=True)):
        if has_header and i == 0:
            continue
        value = row[0]
        if value is None:
            continue
        text = str(value).strip()
        if text:
            names.append(text)

    wb.close()
    return names


# ---------------------------------------------------------------------------
# استبدال النص داخل الأشكال مع الحفاظ على التنسيق
# ---------------------------------------------------------------------------
def _replace_in_text_frame(text_frame, placeholder, value):
    """
    يستبدل placeholder بالقيمة داخل إطار نص. المشكلة أن PowerPoint قد يقسّم
    النص إلى عدة runs، لذا نجمع نص الفقرة، نستبدل، ثم نضع الناتج في أول run
    ونفرّغ الباقي حفاظًا على تنسيق الـ run الأول.
    """
    replaced = False
    for paragraph in text_frame.paragraphs:
        runs = paragraph.runs
        if not runs:
            continue
        full_text = "".join(run.text for run in runs)
        if placeholder not in full_text:
            continue
        new_text = full_text.replace(placeholder, value)
        runs[0].text = new_text
        for run in runs[1:]:
            run.text = ""
        replaced = True
    return replaced


def replace_placeholder(slide, placeholder, value):
    """يستبدل المكان المحجوز في كل أشكال الشريحة (بما فيها داخل الجداول)."""
    count = 0
    for shape in slide.shapes:
        if shape.has_text_frame:
            if _replace_in_text_frame(shape.text_frame, placeholder, value):
                count += 1
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    if _replace_in_text_frame(cell.text_frame, placeholder, value):
                        count += 1
    return count


# ---------------------------------------------------------------------------
# نسخ شريحة داخل نفس العرض التقديمي (مع نسخ العلاقات/الصور)
# ---------------------------------------------------------------------------
def duplicate_slide(prs, source_slide):
    """ينسخ شريحة موجودة ويضيفها في نهاية العرض، مع نسخ العلاقات (صور..)."""
    # استخدام نفس التخطيط (layout) الخاص بالشريحة المصدر
    new_slide = prs.slides.add_slide(source_slide.slide_layout)

    # إزالة أي أشكال أضافها التخطيط تلقائيًا كي ننسخ من المصدر فقط
    for shape in list(new_slide.shapes):
        shape._element.getparent().remove(shape._element)

    # نسخ كل أشكال الشريحة المصدر
    for shape in source_slide.shapes:
        new_el = copy.deepcopy(shape._element)
        new_slide.shapes._spTree.append(new_el)

    # نسخ العلاقات (الصور والوسائط) حتى لا تنكسر بعد النسخ
    for rel_id, rel in source_slide.part.rels.items():
        if "notesSlide" in rel.reltype:
            continue
        if rel.is_external:
            new_slide.part.rels.get_or_add_ext_rel(rel.reltype, rel.target_ref)
        else:
            new_slide.part.relate_to(rel.target_part, rel.reltype, rel_id)

    return new_slide


def delete_slide(prs, slide):
    """يحذف شريحة من العرض التقديمي (يزيل الشريحة وعلاقتها)."""
    xml_slides = prs.slides._sldIdLst
    slide_ids = list(xml_slides)
    rel_ns = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
    for idx, s in enumerate(prs.slides):
        if s == slide:
            rId = slide_ids[idx].get(rel_ns)
            if rId:
                prs.part.drop_rel(rId)
            xml_slides.remove(slide_ids[idx])
            return


# ---------------------------------------------------------------------------
# البناء الرئيسي
# ---------------------------------------------------------------------------
def build(pptx_path, names, placeholder, slide_number, output_path, split=False):
    if not os.path.exists(pptx_path):
        sys.exit(f"خطأ: ملف PowerPoint غير موجود: {pptx_path}")
    if not names:
        sys.exit("خطأ: لم يتم العثور على أي أسماء في ملف Excel.")

    slide_index = slide_number - 1

    if split:
        # ملف مستقل لكل اسم
        os.makedirs(output_path, exist_ok=True)
        generated = []
        for name in names:
            prs = Presentation(pptx_path)
            if slide_index >= len(prs.slides):
                sys.exit(f"خطأ: رقم الشريحة {slide_number} خارج نطاق القالب.")
            replace_placeholder(prs.slides[slide_index], placeholder, name)
            safe = "".join(c if c.isalnum() or c in " _-" else "_" for c in name).strip()
            out_file = os.path.join(output_path, f"{safe or 'name'}.pptx")
            prs.save(out_file)
            generated.append(out_file)
        return generated

    # ملف واحد: شريحة لكل اسم
    prs = Presentation(pptx_path)
    if slide_index >= len(prs.slides):
        sys.exit(f"خطأ: رقم الشريحة {slide_number} خارج نطاق القالب.")

    template_slide = prs.slides[slide_index]

    # ننسخ القالب لكل اسم (مع إبقاء القالب سليمًا فيه المكان المحجوز)،
    # ثم نستبدل الاسم في كل نسخة، وأخيرًا نحذف شريحة القالب الأصلية.
    for name in names:
        dup = duplicate_slide(prs, template_slide)
        replace_placeholder(dup, placeholder, name)

    delete_slide(prs, template_slide)

    prs.save(output_path)
    return [output_path]


def main():
    parser = argparse.ArgumentParser(
        description="إدراج الأسماء من ملف Excel داخل قالب PowerPoint.",
        formatter_class=argparse.RawDescriptionHelpFormatter,
    )
    parser.add_argument("--pptx", required=True, help="مسار قالب PowerPoint")
    parser.add_argument("--excel", required=True, help="مسار ملف Excel الذي يحوي الأسماء")
    parser.add_argument("--output", default="result.pptx",
                        help="مسار ملف الناتج (أو مجلد الناتج مع --split)")
    parser.add_argument("--placeholder", default="{{name}}",
                        help="النص المحجوز داخل القالب (الافتراضي {{name}})")
    parser.add_argument("--column", default=None,
                        help="عمود الأسماء في Excel: حرف مثل A أو رقم مثل 1 (الافتراضي أول عمود)")
    parser.add_argument("--sheet", default=None, help="اسم الورقة (الافتراضي أول ورقة)")
    parser.add_argument("--has-header", action="store_true",
                        help="تجاهُل الصف الأول باعتباره عنوان العمود")
    parser.add_argument("--slide", type=int, default=1,
                        help="رقم شريحة القالب التي تُنسخ (الافتراضي 1)")
    parser.add_argument("--split", action="store_true",
                        help="إنشاء ملف مستقل لكل اسم بدلًا من ملف واحد")

    args = parser.parse_args()

    names = read_names(args.excel, column=args.column, sheet=args.sheet,
                       has_header=args.has_header)
    print(f"تم قراءة {len(names)} اسم من: {args.excel}")

    outputs = build(args.pptx, names, args.placeholder, args.slide,
                    args.output, split=args.split)

    if args.split:
        print(f"تم إنشاء {len(outputs)} ملف داخل المجلد: {args.output}")
    else:
        print(f"تم إنشاء الملف: {outputs[0]}  (عدد الشرائح = {len(names)})")


if __name__ == "__main__":
    main()
